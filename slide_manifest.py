"""
Reads a .pptx file and extracts per-slide text into a structured manifest.
Run once before presenting:

    python slide_manifest.py path/to/deck.pptx

Writes manifest.json to the same directory as the deck. The semantic detector
loads this at startup to know what each slide contains.
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


def extract_manifest(pptx_path: str) -> list[dict]:
    """
    Parse a .pptx file and return a list of slide dicts:
      {
        "slide": 1,          # 1-indexed
        "title": "...",      # text of the title placeholder, or ""
        "body": "...",       # all other text on the slide, joined by spaces
        "notes": "...",      # speaker notes, if any
        "summary": "..."     # title + body combined — used as LLM context
      }
    """
    prs = Presentation(pptx_path)
    manifest = []

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts = []
        notes = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = " ".join(
                para.text.strip()
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            )
            if not text:
                continue

            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0:      # idx 0 = title
                    title = text
                    continue
                # all other placeholder types fall through to body

            body_parts.append(text)

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = " ".join(
                para.text.strip()
                for para in notes_frame.paragraphs
                if para.text.strip()
            )

        body = " ".join(body_parts).strip()
        summary = f"{title}. {body}".strip(". ") if body else title

        manifest.append({
            "slide": i,
            "title": title,
            "body": body,
            "notes": notes,
            "summary": summary,
        })

    return manifest


def save_manifest(manifest: list[dict], output_path: str) -> None:
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2, ensure_ascii=False)
    print(f"[Manifest] Saved {len(manifest)} slides → {output_path}")


def load_manifest(manifest_path: str) -> list[dict]:
    """Load a previously saved manifest JSON. Used by semantic_detector.py."""
    with open(manifest_path, "r", encoding="utf-8") as f:
        return json.load(f)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python slide_manifest.py path/to/deck.pptx")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = str(Path(pptx_path).with_suffix(".manifest.json"))

    print(f"[Manifest] Reading {pptx_path}...")
    manifest = extract_manifest(pptx_path)
    save_manifest(manifest, output_path)

    # Print a preview
    print("\nPreview (first 3 slides):")
    for slide in manifest[:3]:
        print(f"  Slide {slide['slide']}: {slide['title']!r}")
        if slide['body']:
            preview = slide['body'][:80] + ("..." if len(slide['body']) > 80 else "")
            print(f"    Body: {preview}")
